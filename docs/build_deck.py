"""Build the presentation deck for the Local Data Sanitization Pipeline.

Run: python docs/build_deck.py
Outputs: docs/Sanitization_Pipeline_Presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# Palette
NAVY        = RGBColor(0x0F, 0x17, 0x2A)
SLATE       = RGBColor(0x47, 0x55, 0x69)
SLATE_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
SLATE_BG    = RGBColor(0xF1, 0xF5, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
BLUE_SOFT   = RGBColor(0xDB, 0xEA, 0xFE)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT  = RGBColor(0xFE, 0xF3, 0xC7)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_SOFT  = RGBColor(0xDC, 0xFC, 0xE7)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_SOFT    = RGBColor(0xFE, 0xE2, 0xE2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"


# ---- helpers ---------------------------------------------------------------

def add_text(slide, left, top, width, height, text, *,
             font=FONT_BODY, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top  = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=NAVY, bullet_color=BLUE,
                line_spacing=1.2, space_after=4, font=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        rd = p.add_run()
        rd.text = "•  "
        rd.font.name = font
        rd.font.size = Pt(size)
        rd.font.bold = True
        rd.font.color.rgb = bullet_color
        rt = p.add_run()
        rt.text = item
        rt.font.name = font
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *, fill=WHITE, line=None,
             line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.text_frame.margin_left = Pt(0); s.text_frame.margin_right = Pt(0)
    s.text_frame.margin_top = Pt(0); s.text_frame.margin_bottom = Pt(0)
    return s


def add_label_in_shape(shape, text, *, size=12, bold=False, color=NAVY,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                       font=FONT_BODY):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_arrow(slide, x1, y1, x2, y2, *, color=SLATE, width_pt=2.0):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    ln = line.line._get_or_add_ln()
    etree.SubElement(ln, qn("a:headEnd"))
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return line


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_chrome(slide, *, slide_num, total):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=BLUE)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(7.5), Inches(0.3),
             "Local Data Sanitization Pipeline", size=10, color=SLATE_LIGHT)
    add_text(slide, Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3),
             f"{slide_num} / {total}", size=10, color=SLATE_LIGHT,
             align=PP_ALIGN.RIGHT)


def add_title(slide, title):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=NAVY, font=FONT_HEAD)
    add_rect(slide, Inches(0.6), Inches(1.15), Inches(0.6), Inches(0.05),
             fill=AMBER)


# ---- slides ----------------------------------------------------------------

def slide_cover(prs, _n=None, _total=None):
    s = slide_blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, Inches(0), Inches(6.0), SLIDE_W, Inches(0.06), fill=AMBER)

    add_text(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.3),
             "Local Data Sanitization Pipeline",
             size=44, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.6),
             "De-identifying enterprise exports, with a row-level audit trail.",
             size=18, color=SLATE_LIGHT, font=FONT_HEAD)

    add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4),
             "data_engineer_assesment    ·    47 tests    ·    stdlib only",
             size=12, color=AMBER, font=FONT_MONO)


def slide_overview(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "What it does")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.55),
             "Walks a folder of mixed enterprise exports, replaces configured PII "
             "with stable tokens, masks unknown emails / phones, and writes "
             "a sanitized output tree plus six audit reports.",
             size=15, color=NAVY)

    add_text(s, Inches(0.6), Inches(2.55), Inches(6), Inches(0.4),
             "Run", size=12, bold=True, color=BLUE)
    add_rect(s, Inches(0.6), Inches(2.95), Inches(7.5), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.8), Inches(2.96), Inches(7.3), Inches(0.5),
             "python -m sanitizer --input sample_input --output output",
             size=13, color=AMBER, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(3.85), Inches(6), Inches(0.4),
             "Produces", size=12, bold=True, color=BLUE)
    tree = (
        "output/\n"
        "  sanitized/                   <mirrored input tree>\n"
        "  reports/\n"
        "    run_summary.json\n"
        "    file_manifest.jsonl\n"
        "    validation_report.json\n"
        "    pii_transformations.csv\n"
        "    pii_quarantine.csv\n"
        "    analytics.html"
    )
    add_rect(s, Inches(0.6), Inches(4.25), Inches(7.5), Inches(2.7),
             fill=SLATE_BG, line=SLATE_LIGHT)
    add_text(s, Inches(0.8), Inches(4.3), Inches(7.3), Inches(2.6),
             tree, size=12, color=NAVY, font=FONT_MONO)

    # Right column: stats
    stats = [
        ("4",  "supported formats", ".txt / .md / .json / .csv"),
        ("6",  "audit artifacts",   "five JSON/CSV + one HTML"),
        ("4",  "validation checks", "run after artifacts are written"),
        ("47", "tests",             "~1 second wall time"),
        ("0",  "runtime deps",      "stdlib only; pytest is dev-only"),
    ]
    y = Inches(2.55)
    for n_, label, sub in stats:
        add_rect(s, Inches(8.4), y, Inches(4.3), Inches(0.78),
                 fill=WHITE, line=SLATE_LIGHT)
        add_text(s, Inches(8.55), y + Inches(0.12), Inches(0.9), Inches(0.6),
                 n_, size=22, bold=True, color=AMBER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.5), y + Inches(0.06), Inches(3.1), Inches(0.4),
                 label, size=13, bold=True, color=NAVY)
        add_text(s, Inches(9.5), y + Inches(0.4), Inches(3.1), Inches(0.4),
                 sub, size=10, color=SLATE)
        y += Inches(0.88)


def slide_flow(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Flow")

    # Row 1: walker -> classify -> empty? -> processor -> deid
    inp = add_rect(s, Inches(0.5), Inches(2.2), Inches(1.5), Inches(0.9),
                   fill=NAVY, shape=MSO_SHAPE.CAN)
    add_label_in_shape(inp, "input/", size=12, bold=True, color=WHITE)

    walker = add_rect(s, Inches(2.3), Inches(2.2), Inches(1.7), Inches(0.9),
                      fill=BLUE_SOFT, line=BLUE)
    add_label_in_shape(walker, "walk\nsorted", size=11, bold=True, color=NAVY)

    d1 = add_rect(s, Inches(4.3), Inches(2.1), Inches(1.4), Inches(1.1),
                  fill=AMBER_SOFT, line=AMBER, shape=MSO_SHAPE.DIAMOND)
    add_label_in_shape(d1, "supported?", size=10, bold=True, color=NAVY)

    d2 = add_rect(s, Inches(6.1), Inches(2.1), Inches(1.4), Inches(1.1),
                  fill=AMBER_SOFT, line=AMBER, shape=MSO_SHAPE.DIAMOND)
    add_label_in_shape(d2, "empty?", size=10, bold=True, color=NAVY)

    proc = add_rect(s, Inches(7.9), Inches(2.2), Inches(1.7), Inches(0.9),
                    fill=BLUE_SOFT, line=BLUE)
    add_label_in_shape(proc, "processor\ntxt/md/json/csv", size=10,
                       bold=True, color=NAVY)

    deid = add_rect(s, Inches(10.0), Inches(2.2), Inches(2.8), Inches(0.9),
                    fill=AMBER, line=AMBER)
    add_label_in_shape(deid, "DeIdentifier.apply", size=12,
                       bold=True, color=NAVY)

    # Branch boxes
    sk = add_rect(s, Inches(4.3), Inches(0.95), Inches(1.4), Inches(0.6),
                  fill=SLATE_BG, line=SLATE_LIGHT)
    add_label_in_shape(sk, "skipped_unsupported", size=9, color=SLATE)

    em = add_rect(s, Inches(6.1), Inches(0.95), Inches(1.4), Inches(0.6),
                  fill=SLATE_BG, line=SLATE_LIGHT)
    add_label_in_shape(em, "empty\n(0-byte mirror)", size=9, color=SLATE)

    fl = add_rect(s, Inches(8.5), Inches(3.6), Inches(1.4), Inches(0.6),
                  fill=RED_SOFT, line=RED)
    add_label_in_shape(fl, "failed\n(captured)", size=9, color=NAVY)

    # Arrows row 1
    add_arrow(s, Inches(2.0), Inches(2.65), Inches(2.3), Inches(2.65))
    add_arrow(s, Inches(4.0), Inches(2.65), Inches(4.3), Inches(2.65))
    add_arrow(s, Inches(5.7), Inches(2.65), Inches(6.1), Inches(2.65))
    add_arrow(s, Inches(7.5), Inches(2.65), Inches(7.9), Inches(2.65))
    add_arrow(s, Inches(9.6), Inches(2.65), Inches(10.0), Inches(2.65))
    add_arrow(s, Inches(5.0), Inches(2.1), Inches(5.0), Inches(1.55))
    add_arrow(s, Inches(6.8), Inches(2.1), Inches(6.8), Inches(1.55))
    add_arrow(s, Inches(8.75), Inches(3.1), Inches(8.75), Inches(3.6))

    # Row 2: outputs
    out_y = Inches(4.7)
    outs = [
        ("sanitized/\n(mirrored)",    GREEN_SOFT, GREEN),
        ("file_manifest.jsonl",        GREEN_SOFT, GREEN),
        ("pii_transformations.csv",    GREEN_SOFT, GREEN),
        ("pii_quarantine.csv",         GREEN_SOFT, GREEN),
        ("validation_report.json",     AMBER_SOFT, AMBER),
        ("run_summary.json",           AMBER_SOFT, AMBER),
        ("analytics.html",             BLUE_SOFT,  BLUE),
    ]
    box_w = Inches(1.7)
    pad = Inches(0.05)
    x = Inches(0.5)
    for label, fill, accent in outs:
        b = add_rect(s, x, out_y, box_w, Inches(0.85), fill=fill, line=accent)
        add_label_in_shape(b, label, size=10, bold=True, color=NAVY)
        x += box_w + pad

    # Down arrow from deid to outputs row
    add_arrow(s, Inches(11.0), Inches(3.1), Inches(11.0), Inches(4.7))

    # Note row
    add_text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.45),
             "The validator re-scans the written sanitized files and "
             "cross-checks the manifest against the input tree.",
             size=12, color=SLATE)
    add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.45),
             "Per-file try/except: a malformed file becomes a manifest row, not a crashed run.",
             size=12, color=SLATE)


def slide_supported(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Supported file types")

    headers = ["extension", "reader", "writer", "location format"]
    rows = [
        (".txt / .md", "UTF-8 text",       "UTF-8 text",                    'line N, column M'),
        (".json",      "json.load",        "json.dumps(indent=2)",          '$.issues[1].comments[0].body'),
        (".csv",       "csv.DictReader",   "csv.DictWriter",                'row N, column "from"'),
    ]
    col_xs = [Inches(0.6), Inches(2.8), Inches(5.0), Inches(7.5)]
    col_ws = [Inches(2.2), Inches(2.2), Inches(2.5), Inches(5.2)]
    y = Inches(1.6)
    for x, w, h in zip(col_xs, col_ws, headers):
        rh = add_rect(s, x, y, w, Inches(0.5), fill=NAVY)
        add_label_in_shape(rh, h, size=11, bold=True, color=AMBER,
                           align=PP_ALIGN.LEFT)
        rh.text_frame.margin_left = Pt(10)
    y += Inches(0.5)
    alt = False
    for row in rows:
        bg = SLATE_BG if alt else WHITE
        for x, w, val in zip(col_xs, col_ws, row):
            cell = add_rect(s, x, y, w, Inches(0.6), fill=bg,
                            line=SLATE_LIGHT, line_w=0.5)
            add_label_in_shape(cell, val, size=12, color=NAVY,
                               font=FONT_MONO, align=PP_ALIGN.LEFT)
            cell.text_frame.margin_left = Pt(10)
        y += Inches(0.6)
        alt = not alt

    add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
             "Anything else", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(5.05), Inches(12), Inches(1.5),
             ".pdf, .png, .xlsx, .zip, ... are recorded in the manifest "
             "as skipped_unsupported. Their bytes are never opened.",
             size=13, color=SLATE)


def slide_replacement_order(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Replacement order")

    add_rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.65),
             fill=NAVY)
    add_text(s, Inches(0.85), Inches(1.6), Inches(11.7), Inches(0.55),
             "emails  →  phones  →  persons (longest first)  →  organizations (longest first)",
             size=18, bold=True, color=AMBER, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Wrong vs right
    add_text(s, Inches(0.6), Inches(2.55), Inches(6), Inches(0.4),
             "Orgs first (wrong)", size=14, bold=True, color=RED)
    add_rect(s, Inches(0.6), Inches(2.95), Inches(6.0), Inches(2.5),
             fill=RED_SOFT, line=RED)
    wrong = [
        "input:    sarah@betahealth.io",
        "org:      BetaHealth → ORG_001",
        "result:   sarah@ORG_001.io",
        "email:    no match",
        "leak:     'sarah' + '.io' raw",
    ]
    y = Inches(3.07)
    for line in wrong:
        add_text(s, Inches(0.85), y, Inches(5.6), Inches(0.4), line,
                 size=12, color=NAVY, font=FONT_MONO)
        y += Inches(0.46)

    add_text(s, Inches(6.9), Inches(2.55), Inches(6), Inches(0.4),
             "Emails first (correct)", size=14, bold=True, color=GREEN)
    add_rect(s, Inches(6.9), Inches(2.95), Inches(6.0), Inches(2.5),
             fill=GREEN_SOFT, line=GREEN)
    right = [
        "input:    sarah@betahealth.io",
        "email:    whole address → EMAIL_002",
        "result:   EMAIL_002",
        "org:      nothing left to match",
        "clean:    no leakage",
    ]
    y = Inches(3.07)
    for line in right:
        add_text(s, Inches(7.15), y, Inches(5.6), Inches(0.4), line,
                 size=12, color=NAVY, font=FONT_MONO)
        y += Inches(0.46)

    add_text(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.45),
             "Longest-alias-first applies the same logic within a group: "
             "'Sarah Chen' before 'Sarah' so we don't leave 'PERSON_002 Chen'.",
             size=12, color=SLATE)


def slide_three_phase(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Three-phase apply()")

    phases = [
        ("Phase 1",
         "Scan original text",
         "Find all match offsets in the original string. "
         "Higher-priority matches register their span; later "
         "matches inside a covered span are dropped.",
         BLUE),
        ("Phase 2",
         "Cascading replacement",
         "Run emails → phones → persons → orgs in order. "
         "Mapped values become tokens; unmapped emails / phones "
         "become <UNMAPPED_*> placeholders.",
         AMBER),
        ("Phase 3",
         "Render snippets",
         "Render snippet windows from the fully sanitized text — "
         "so the focal value and any neighboring PII appear as tokens.",
         GREEN),
    ]
    x = Inches(0.6)
    for label, title, body, accent in phases:
        add_rect(s, x, Inches(1.55), Inches(4.0), Inches(4.4),
                 fill=WHITE, line=accent, line_w=1.5)
        add_rect(s, x, Inches(1.55), Inches(4.0), Inches(0.5), fill=accent)
        add_text(s, x, Inches(1.6), Inches(4.0), Inches(0.4),
                 label, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), Inches(2.2), Inches(3.6), Inches(0.5),
                 title, size=17, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), Inches(2.85), Inches(3.6), Inches(2.8),
                 body, size=13, color=SLATE)
        x += Inches(4.2)

    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             "Why split the work: offsets stay tied to the input file the "
             "operator opens, and snippets can never carry a neighbour's raw value.",
             size=12, color=SLATE)


def slide_config(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Configuration: config/entities.json")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
             "Person-centric. Aliases, emails, and phones live under "
             "the entity that owns them.",
             size=13, color=SLATE)

    add_rect(s, Inches(0.6), Inches(2.1), Inches(7.5), Inches(4.7),
             fill=NAVY)
    json_text = (
        '{\n'
        '  "persons": [\n'
        '    {\n'
        '      "canonical_id": "PERSON_002",\n'
        '      "aliases": ["Sarah Chen", "Sarah"],\n'
        '      "emails": [\n'
        '        {"value": "sarah@betahealth.io",\n'
        '         "token": "EMAIL_002"}\n'
        '      ],\n'
        '      "phones": [\n'
        '        {"value": "+1-212-555-0199",\n'
        '         "token": "PHONE_001"}\n'
        '      ]\n'
        '    }\n'
        '  ],\n'
        '  "organizations": [\n'
        '    {"canonical_id": "ORG_002",\n'
        '     "aliases": ["Acme Inc.", "Acme"]}\n'
        '  ]\n'
        '}'
    )
    add_text(s, Inches(0.85), Inches(2.25), Inches(7.2), Inches(4.5),
             json_text, size=12, color=WHITE, font=FONT_MONO)

    add_text(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(0.4),
             "Loader contract", size=14, bold=True, color=BLUE)
    add_bullets(s, Inches(8.4), Inches(2.55), Inches(4.5), Inches(2.0),
                ["Phone keys normalized to digits",
                 "Email keys lowercased",
                 "Raises on missing fields",
                 "Raises on conflicting tokens",
                 "Allows idempotent restatements"],
                size=12, color=NAVY, line_spacing=1.2)

    add_text(s, Inches(8.4), Inches(4.6), Inches(4.5), Inches(0.4),
             "Runtime", size=14, bold=True, color=AMBER)
    add_bullets(s, Inches(8.4), Inches(5.05), Inches(4.5), Inches(2.0),
                ["Flat lookup tables derived at load time",
                 "Alias rules sorted longest-first",
                 "Lookarounds (?<!\\w) ... (?!\\w) instead of \\b"],
                size=12, color=NAVY, line_spacing=1.2)


def slide_mapped_unmapped(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Mapped vs unmapped")

    headers = ["kind", "configured?", "sanitized output", "report"]
    rows = [
        ("email",  "yes", "EMAIL_001",        "pii_transformations.csv"),
        ("phone",  "yes", "PHONE_001",        "pii_transformations.csv"),
        ("person", "yes", "PERSON_001",       "pii_transformations.csv"),
        ("org",    "yes", "ORG_001",          "pii_transformations.csv"),
        ("email",  "no",  "<UNMAPPED_EMAIL>", "pii_quarantine.csv"),
        ("phone",  "no",  "<UNMAPPED_PHONE>", "pii_quarantine.csv"),
    ]
    col_xs = [Inches(0.6), Inches(2.5), Inches(4.7), Inches(8.0)]
    col_ws = [Inches(1.9), Inches(2.2), Inches(3.3), Inches(4.7)]
    y = Inches(1.55)
    for x, w, h in zip(col_xs, col_ws, headers):
        rh = add_rect(s, x, y, w, Inches(0.45), fill=NAVY)
        add_label_in_shape(rh, h, size=11, bold=True, color=AMBER,
                           align=PP_ALIGN.LEFT)
        rh.text_frame.margin_left = Pt(8)
    y += Inches(0.45)
    for row in rows:
        bg = GREEN_SOFT if row[1] == "yes" else AMBER_SOFT
        for x, w, val in zip(col_xs, col_ws, row):
            cell = add_rect(s, x, y, w, Inches(0.45), fill=bg,
                            line=SLATE_LIGHT, line_w=0.5)
            font = (FONT_MONO if val.startswith("<") or val.endswith(".csv")
                    or "_00" in val else FONT_BODY)
            add_label_in_shape(cell, val, size=11, color=NAVY,
                               font=font, align=PP_ALIGN.LEFT)
            cell.text_frame.margin_left = Pt(8)
        y += Inches(0.45)

    add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
             "Why a generic placeholder, not an auto-pseudonym",
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(1.4),
             "An auto-pseudonym like EMAIL_AUTO_a3f9b21c is visually "
             "indistinguishable from a reviewed token — a reader can't "
             "tell whether it's been approved. The placeholder makes the "
             "unreviewed state obvious in the output. Distinctness is "
             "preserved in the quarantine CSV via value_hash.",
             size=13, color=SLATE)


def slide_triage(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Quarantine triage loop")

    steps = [
        ("1", "Detect",        "regex hits a value\nnot in config"),
        ("2", "Mask",          "<UNMAPPED_EMAIL>\n<UNMAPPED_PHONE>"),
        ("3", "Quarantine",    "row in\npii_quarantine.csv"),
        ("4", "Review",        "operator reads\nvalue + snippet"),
        ("5", "Update",        "add to\nentities.json"),
        ("6", "Re-run",        "row migrates to\ntransformations.csv"),
    ]
    x = Inches(0.55)
    box_w = Inches(2.0)
    box_h = Inches(2.4)
    y = Inches(2.2)
    for i, (num, title, body) in enumerate(steps):
        accent = BLUE if i in (0, 3) else (AMBER if i in (1, 2) else GREEN)
        circ = add_rect(s, x + Inches(0.7), y - Inches(0.4),
                        Inches(0.6), Inches(0.6),
                        fill=accent, shape=MSO_SHAPE.OVAL)
        add_label_in_shape(circ, num, size=18, bold=True, color=WHITE)
        add_rect(s, x, y, box_w, box_h, fill=WHITE, line=accent, line_w=1.5)
        add_text(s, x, y + Inches(0.35), box_w, Inches(0.5),
                 title, size=15, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.9), box_w, Inches(1.5),
                 body, size=12, color=SLATE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(s, x + box_w, y + Inches(1.2),
                      x + box_w + Inches(0.1), y + Inches(1.2),
                      color=SLATE, width_pt=2)
        x += Inches(2.13)

    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
             "value_hash = sha256(value)[:8]. Same hash everywhere — "
             "one config update fixes every occurrence across files.",
             size=13, color=SLATE)


def slide_outputs(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Six audit artifacts")

    artifacts = [
        ("run_summary.json",
         "Top-level totals, by-extension histogram, validation summary.\n"
         "One object per run. Designed for dashboards / DAGs.",
         BLUE),
        ("file_manifest.jsonl",
         "One line per input file, including skipped / failed / empty.\n"
         "Input hash, output hash, replacement counts, error.",
         BLUE),
        ("validation_report.json",
         "Four post-run checks against the written artifacts.",
         AMBER),
        ("pii_transformations.csv",
         "One row per mapped replacement. file + kind + value +\n"
         "value_hash + token + status + location + safe snippet.",
         GREEN),
        ("pii_quarantine.csv",
         "Same schema, status='unmapped'. The operator backlog.\n"
         "Only written if at least one row exists.",
         AMBER),
        ("analytics.html",
         "Self-contained dashboard: stat tiles, entity↔file graph,\n"
         "grouped quarantine triage panel.",
         BLUE),
    ]
    cell_w = Inches(6.0)
    cell_h = Inches(1.6)
    pad = Inches(0.15)
    for i, (name, body, accent) in enumerate(artifacts):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * (cell_w + pad)
        y = Inches(1.55) + row * (cell_h + pad)
        add_rect(s, x, y, cell_w, cell_h, fill=WHITE, line=SLATE_LIGHT)
        add_rect(s, x, y, Inches(0.12), cell_h, fill=accent)
        add_text(s, x + Inches(0.3), y + Inches(0.12),
                 cell_w - Inches(0.3), Inches(0.4),
                 name, size=14, bold=True, color=NAVY, font=FONT_MONO)
        add_text(s, x + Inches(0.3), y + Inches(0.55),
                 cell_w - Inches(0.3), cell_h - Inches(0.5),
                 body, size=11, color=SLATE)


def slide_pii_schema(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "PII row schema (shared by both reports)")

    add_rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.55),
             fill=NAVY)
    add_text(s, Inches(0.85), Inches(1.55), Inches(11.7), Inches(0.55),
             "file, kind, value, value_hash, token, status, location, snippet",
             size=14, bold=True, color=AMBER, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cols_info = [
        ("file",        "input file relative path (POSIX)"),
        ("kind",        "email / phone / person / organization"),
        ("value",       "normalized raw value"),
        ("value_hash",  "first 8 chars of sha256(value), cross-doc dedup key"),
        ("token",       "EMAIL_001, PERSON_002, <UNMAPPED_EMAIL>, ..."),
        ("status",      "mapped or unmapped"),
        ("location",    'line N, column M  /  $.path[i].field  /  row N, column "X"'),
        ("snippet",     "~60 chars context, rendered against sanitized text"),
    ]
    y = Inches(2.4)
    for name, desc in cols_info:
        add_rect(s, Inches(0.6), y, Inches(2.4), Inches(0.45),
                 fill=BLUE_SOFT, line=BLUE)
        add_text(s, Inches(0.7), y + Inches(0.05), Inches(2.3), Inches(0.4),
                 name, size=12, bold=True, color=NAVY,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.15), y + Inches(0.05),
                 Inches(9.5), Inches(0.4),
                 desc, size=12, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.5)

    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "A row migrates from quarantine to transformations on the "
             "next run with no other field changes.",
             size=12, color=SLATE)


def slide_validation(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Validation: 4 checks against the written artifacts")

    checks = [
        ("1", "no_raw_emails_in_sanitized_outputs",
         "Re-scan every file under sanitized/ for raw email patterns."),
        ("2", "no_raw_phone_numbers_in_sanitized_outputs",
         "Same, for phone numbers."),
        ("3", "processed_files_have_outputs",
         "Every manifest row with status=processed has an output file on disk."),
        ("4", "all_input_files_accounted_for",
         "Every input file appears exactly once in the manifest."),
    ]
    y = Inches(1.9)
    for num, name, desc in checks:
        circ = add_rect(s, Inches(0.7), y, Inches(0.6), Inches(0.6),
                        fill=GREEN, shape=MSO_SHAPE.OVAL)
        add_label_in_shape(circ, num, size=18, bold=True, color=WHITE)
        add_rect(s, Inches(1.5), y, Inches(11.2), Inches(0.6),
                 fill=GREEN_SOFT, line=GREEN)
        add_text(s, Inches(1.7), y + Inches(0.05), Inches(6.5),
                 Inches(0.5),
                 name, size=12, bold=True, color=NAVY,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.0), y + Inches(0.05),
                 Inches(4.7), Inches(0.5),
                 desc, size=11, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.85)

    add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
             "The validator runs after every other artifact is written, "
             "so the manifest is the input it audits — not in-memory state.",
             size=13, color=SLATE)
    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             "Exit codes: 0 = passed, 2 = warnings or leak, 1 = catastrophic failure.",
             size=13, color=SLATE)


def slide_dashboard(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "analytics.html")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.55),
             "Single-page HTML dashboard rendered per run. Self-contained "
             "except for one CDN script tag (vis-network).",
             size=14, color=SLATE)

    regions = [
        ("Stat tiles",
         "Run id, timestamps, 8 color-coded tiles\n"
         "(discovered / processed / skipped / failed /\n"
         "empty / mapped / unmapped / validation).",
         BLUE),
        ("Entity ↔ file graph",
         "Files as boxes sized by replacement count.\n"
         "Entities as ellipses colored by kind.\n"
         "Edges weighted by occurrence count.",
         AMBER),
        ("Quarantine panel",
         "Unmapped values grouped by\n"
         "(kind, value, hash). Each group lists\n"
         "occurrences with file, location, snippet.",
         GREEN),
    ]
    x = Inches(0.6)
    for title, body, accent in regions:
        add_rect(s, x, Inches(2.4), Inches(4.0), Inches(3.6),
                 fill=WHITE, line=accent, line_w=2)
        add_rect(s, x, Inches(2.4), Inches(4.0), Inches(0.5), fill=accent)
        add_text(s, x, Inches(2.45), Inches(4.0), Inches(0.45),
                 title, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), Inches(3.05),
                 Inches(3.6), Inches(2.9),
                 body, size=12, color=NAVY)
        x += Inches(4.2)

    add_text(s, Inches(0.6), Inches(6.25), Inches(12), Inches(0.5),
             "Header carries back-links to all five sibling artifacts.",
             size=13, color=SLATE)


def slide_determinism(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Determinism")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
             "Sanitized output is byte-identical across reruns on the same "
             "input + config. Five mechanisms:",
             size=14, color=SLATE)

    mechs = [
        ("Sorted traversal",
         "dirnames.sort() + sorted(filenames) at every os.walk level"),
        ("Deterministic placeholders",
         "<UNMAPPED_EMAIL> / <UNMAPPED_PHONE> — no per-run randomness"),
        ("Content-derived value_hash",
         "sha256(value)[:8] — pure function of the input value"),
        ("Stable column order + LF",
         "hard-coded PII_FIELDNAMES, lineterminator='\\n'"),
        ("Deterministic JSON",
         "indent=2, dict insertion order preserved (Python 3.7+)"),
    ]
    y = Inches(2.3)
    for title, desc in mechs:
        add_rect(s, Inches(0.6), y, Inches(4.0), Inches(0.55),
                 fill=BLUE_SOFT, line=BLUE)
        add_text(s, Inches(0.8), y + Inches(0.05), Inches(3.8), Inches(0.5),
                 title, size=13, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(4.7), y, Inches(8.0), Inches(0.55),
                 fill=WHITE, line=SLATE_LIGHT)
        add_text(s, Inches(4.85), y + Inches(0.05),
                 Inches(7.8), Inches(0.5),
                 desc, size=12, color=SLATE,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.65)

    add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
             "Only run_id and started_at / completed_at differ between runs — by design.",
             size=13, color=SLATE)


def slide_testing(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Testing")

    add_rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.7),
             fill=NAVY)
    add_text(s, Inches(0.85), Inches(1.6), Inches(11.7), Inches(0.6),
             "47 tests  ·  ~1 second wall time  ·  "
             "tests/test_deid.py + test_pipeline.py + test_validation.py",
             size=14, bold=True, color=AMBER,
             font=FONT_MONO, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.4),
             "What the suite protects",
             size=14, bold=True, color=BLUE)
    items = [
        "Replacement-order invariants (email-before-org, longest-alias-first)",
        "Regex boundary cases (ISO timestamps not matching phones, aliases not matching inside other words)",
        "Config schema invariants (conflict detection, missing fields, idempotent duplicates)",
        "Both PII reports — location accuracy per processor, snippet privacy, schema sharing, RFC-4180 round-trips",
        "Summary-vs-row totals consistency",
        "Analytics dashboard — file written, embedded JSON shape, snippets don't leak raw PII",
        "Four validation checks under both clean and tampered conditions",
        "Byte-identical sanitized files across repeated runs",
    ]
    add_bullets(s, Inches(0.7), Inches(3.0), Inches(12), Inches(4),
                items, size=13, color=NAVY, bullet_color=GREEN,
                line_spacing=1.25)


def slide_tradeoffs(prs, n, total):
    s = slide_blank(prs)
    add_chrome(s, slide_num=n, total=total)
    add_title(s, "Tradeoffs & limitations")

    rows = [
        ("Stdlib only at runtime",
         "Easy to audit, portable, no supply chain.",
         "No libphonenumber, no streaming JSON, no NER libraries."),
        ("Regex + explicit config",
         "Deterministic and explainable; every replacement is reviewable.",
         "Misses pronouns, unknown names, role references."),
        ("Per-file failure isolation",
         "One bad file doesn't abort the run.",
         "Partial success is only safe if downstream respects warnings."),
        ("<UNMAPPED_*> placeholders",
         "Unreviewed state is visible in the output.",
         "Sanitized text loses unknown-entity distinctness."),
        ("Plain SHA value_hash",
         "Stable cross-doc dedup key.",
         "Guessable. Production wants HMAC with a managed key."),
        ("PII CSVs carry raw values",
         "Operator can act on findings directly.",
         "Reports need stricter access controls than sanitized outputs."),
        ("Single-process runtime",
         "Simple, reproducible, fine for the sample.",
         "Not enough for very large exports — needs streaming + workers."),
        ("No PDF / image / spreadsheet",
         "Supported surface stays explicit and reviewable.",
         "Skipped files may still contain PII."),
    ]
    headers = ["Decision", "What we get", "What we accept"]
    col_xs = [Inches(0.6), Inches(4.5), Inches(8.7)]
    col_ws = [Inches(3.85), Inches(4.15), Inches(4.05)]
    y = Inches(1.55)
    for x, w, h in zip(col_xs, col_ws, headers):
        rh = add_rect(s, x, y, w, Inches(0.45), fill=NAVY)
        add_label_in_shape(rh, h, size=11, bold=True,
                           color=AMBER, align=PP_ALIGN.LEFT)
        rh.text_frame.margin_left = Pt(10)
    y += Inches(0.45)
    for i, row in enumerate(rows):
        bg = SLATE_BG if i % 2 == 0 else WHITE
        h = Inches(0.62)
        for x, w, val in zip(col_xs, col_ws, row):
            cell = add_rect(s, x, y, w, h, fill=bg,
                            line=SLATE_LIGHT, line_w=0.5)
            add_label_in_shape(cell, val, size=11, color=NAVY,
                               align=PP_ALIGN.LEFT,
                               anchor=MSO_ANCHOR.MIDDLE)
            cell.text_frame.margin_left = Pt(10)
            cell.text_frame.margin_right = Pt(8)
        y += h


def slide_qa(prs, _n=None, _total=None):
    s = slide_blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, Inches(0), Inches(6.0), SLIDE_W, Inches(0.06), fill=AMBER)
    add_text(s, Inches(0.8), Inches(2.8), Inches(12), Inches(1.3),
             "Questions?",
             size=60, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(4.2), Inches(12), Inches(0.5),
             "Local Data Sanitization Pipeline",
             size=16, color=SLATE_LIGHT, font=FONT_HEAD)


# ---- build -----------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,                # 1
        slide_overview,             # 2
        slide_flow,                 # 3
        slide_supported,            # 4
        slide_replacement_order,    # 5
        slide_three_phase,          # 6
        slide_config,               # 7
        slide_mapped_unmapped,      # 8
        slide_triage,               # 9
        slide_outputs,              # 10
        slide_pii_schema,           # 11
        slide_validation,           # 12
        slide_dashboard,            # 13
        slide_determinism,          # 14
        slide_testing,              # 15
        slide_tradeoffs,            # 16
        slide_qa,                   # 17
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        fn(prs, i, total)

    out = Path(__file__).parent / "Sanitization_Pipeline_Presentation.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({total} slides)")


if __name__ == "__main__":
    main()
